# Refactor plan & prototype
# 1. Herstructurering context datastructuur naar "variables" + "collections"
# 2. Mapping naar herbruikbare slide rendering

from .pptx_utils import duplicate_slide, get_shape_text, is_placeholder
from .image_utils import replace_image
from pptx import Presentation
from jinja2 import Template as JinjaTemplate
import re

def build_rendered_presentation(template_path, context, env):
    pres = Presentation(template_path)

    collections = context.get("collections", {})
    split_config = context.get("split", {})
    variables = context.get("variables", {})

    slides_to_remove = []

    for idx, slide in enumerate(list(pres.slides)):
        collection_name = detect_repeating_slide(slide)
        if collection_name and collection_name in collections:
            config = split_config.get(collection_name, {"type": "auto"})
            rows = collections[collection_name]["rows"]
            chunks = paginate(rows, config)

            for chunk in chunks:
                new_slide = duplicate_slide(pres, slide)
                # Hier:
                render_slide_with_context(new_slide, {
                    **variables,
                    **collections,
                    "_splitslide_chunk": chunk  # ✅ veilige, tijdelijke naam
                }, env)

            slides_to_remove.append(slide)
        else:
            render_slide_with_context(slide, {**variables, **collections}, env)

    for slide in slides_to_remove:
        rId = pres.slides._sldIdLst[pres.slides.index(slide)].rId
        pres.part.drop_rel(rId)
        pres.slides._sldIdLst.remove(pres.slides._sldIdLst[pres.slides.index(slide)])

    return pres


def paginate(rows, config):
    if config["type"] == "fixed":
        max_per_slide = config.get("max_per_slide", 1)
        return [rows[i:i + max_per_slide] for i in range(0, len(rows), max_per_slide)]
    return [rows]  # auto = alles in 1 groep


def detect_repeating_slide(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = get_shape_text(shape)
            matches = re.findall(r"{{\s*(\w+)\.", text)
            for match in matches:
                return match  # Return the first detected collection name
    return None

def render_slide_with_context(slide, context, env):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = get_shape_text(shape)
            if "{{" in text or "{%" in text:
                try:
                    template = env.from_string(text)
                    rendered = template.render(**context)
                    print("Template text:", text)
                    print(f"Rendered output: {rendered}")
                    shape.text_frame.clear()
                    shape.text_frame.text = rendered
                except Exception as e:
                    print(f"Fout bij renderen van '{text}': {e}")

        elif shape.name.startswith("{{") and shape.name.endswith("}}"):
            key = shape.name.strip("{} ")
            path = context.get(key)
            if path:
                try:
                    replace_image(shape, path)
                except Exception as e:
                    print(f"Fout bij vervangen afbeelding '{key}': {e}")

