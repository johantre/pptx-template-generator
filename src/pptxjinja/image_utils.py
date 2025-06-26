from pptx.slide import Slide
from pptx.parts.image import ImagePart

def replace_image(shape, image_path):
    # Zoek slide
    slide = shape
    while not isinstance(slide, Slide):
        slide = getattr(slide, "_parent", None)
        if slide is None:
            raise RuntimeError("Kon slide niet vinden van shape")

    # Zoek embed-relatie-id
    blip = shape._element.xpath(".//a:blip")
    if not blip:
        rel_id = None
    else:
        rel_id = blip[0].get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")

    # Bepaal of deze relatie exclusief door dit shape wordt gebruikt
    drop_relation = False
    if rel_id:
        count = 0
        for sld in slide.part.package.iter_parts():
            rls = getattr(sld, "rels", {})
            if rel_id in rls._rels:
                count += 1
        if count == 1:
            drop_relation = True

    # Indien exclusief, verwijder de relatie
    if drop_relation:
        slide.part.drop_rel(rel_id)

    # Verwijder shape uit slide
    slide.shapes._spTree.remove(shape._element)

    # Voeg nieuwe afbeelding toe
    slide.shapes.add_picture(image_path, shape.left, shape.top, shape.width, shape.height)
